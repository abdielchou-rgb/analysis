"""V54 505 Toolkit Style DNA Extractor
=========================================
Extracts real institutional chart styling parameters from
investment bank and consulting firm PPT files.

Goal: Replace hand-written institution color palettes with
parameters extracted from actual presentation charts.

Output: data/505_institution_style_dna.json
"""

import json
import logging
from pathlib import Path

logger = logging.getLogger("v54.style_dna")

# Paths
PROJECT_ROOT = Path(__file__).resolve().parent.parent
TOOLKIT_DIR = PROJECT_ROOT / "data" / "505 九大投行咨询公司麦肯锡高盛摩根贝恩罗兰贝格埃森哲德勤JP摩根图表toolkit"
PPT_DIR = TOOLKIT_DIR / "咨询名企原版PPT"
OUTPUT_PATH = PROJECT_ROOT / "data" / "505_institution_style_dna.json"

# Institution name patterns (file name → institution ID)
INSTITUTION_PATTERNS = {
    "goldman": "gs",
    "高盛": "gs",
    "morgan": "ms",
    "摩根": "ms",
    "jpmorgan": "jpm",
    "jp morgan": "jpm",
    "mckinsey": "mck",
    "麦肯锡": "mck",
    "bcg": "bcg",
    "bain": "bain",
    "贝恩": "bain",
    "roland berger": "roland_berger",
    "罗兰贝格": "roland_berger",
    "埃森哲": "accenture",
    "accenture": "accenture",
    "deloitte": "deloitte",
    "德勤": "deloitte",
    "pwc": "pwc",
    "普华永道": "pwc",
    "ey": "ey",
    "安永": "ey",
    "kpmg": "kpmg",
    "毕马威": "kpmg",
    "citi": "citi",
    "花旗": "citi",
    "cicc": "cicc",
    "中金": "cicc",
    "citic": "citic",
    "中信": "citic",
}


def identify_institution(filename: str) -> str | None:
    """Identify institution from filename."""
    name_lower = filename.lower()
    for pattern, inst_id in INSTITUTION_PATTERNS.items():
        if pattern in name_lower:
            return inst_id
    return None


def extract_colors_from_shape(shape) -> list[str]:
    """Try to extract colors from a PPT shape."""
    colors = []
    try:
        # Check fill color
        if hasattr(shape, "fill") and shape.fill:
            if hasattr(shape.fill, "fore_color") and shape.fill.fore_color:
                try:
                    colors.append(str(shape.fill.fore_color.rgb))
                except Exception:
                    pass
        # Check line color
        if hasattr(shape, "line") and shape.line:
            if hasattr(shape.line, "color") and shape.line.color:
                try:
                    colors.append(str(shape.line.color.rgb))
                except Exception:
                    pass
        # Check font color
        if hasattr(shape, "text_frame") and shape.text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.color and run.font.color.rgb:
                        colors.append(str(run.font.color.rgb))
    except Exception:
        pass
    return colors


def extract_chart_colors(ppt_path: Path, chart) -> dict:
    """Extract color palette from a chart object."""
    result = {
        "palette": [],
        "chart_type": None,
        "has_legend": False,
        "has_data_labels": False,
    }
    try:
        if hasattr(chart, "chart_type"):
            result["chart_type"] = str(chart.chart_type)

        # Try to extract series colors
        if hasattr(chart, "series"):
            for series in chart.series:
                try:
                    if hasattr(series, "format") and series.format:
                        if hasattr(series.format, "fill") and series.format.fill:
                            if hasattr(series.format.fill, "fore_color") and series.format.fill.fore_color:
                                try:
                                    color_val = str(series.format.fill.fore_color.rgb)
                                    if color_val not in result["palette"]:
                                        result["palette"].append(color_val)
                                except Exception:
                                    pass
                except Exception:
                    pass

        # Check for legend
        if hasattr(chart, "has_legend"):
            result["has_legend"] = chart.has_legend

        # Check for data labels
        if hasattr(chart, "plot") and chart.plot:
            if hasattr(chart.plot, "has_data_labels"):
                result["has_data_labels"] = chart.plot.has_data_labels
    except Exception:
        pass

    return result


def scan_ppt(ppt_path: Path) -> dict:
    """Scan a single PPT file and extract chart styling parameters."""
    try:
        from pptx import Presentation
        from pptx.util import Emu, Pt  # noqa: F401  (availability probe)
    except ImportError:
        logger.error("python-pptx not installed. Run: pip install python-pptx")
        return {}

    result = {
        "file": ppt_path.name,
        "institution": identify_institution(ppt_path.stem),
        "charts_found": 0,
        "palette": [],
        "font_families": set(),
        "font_sizes": [],
        "grid_style": {},
        "chart_types": {},
        "has_data_labels_total": 0,
        "has_legend_total": 0,
    }

    try:
        prs = Presentation(str(ppt_path))
    except Exception as e:
        logger.warning(f"Cannot open {ppt_path.name}: {e}")
        return result

    total_slides = len(prs.slides)
    shapes_checked = 0

    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            shapes_checked += 1

            # Extract colors from all shapes
            colors = extract_colors_from_shape(shape)
            for c in colors:
                if c not in result["palette"]:
                    result["palette"].append(c)

            # Extract font info
            if hasattr(shape, "text_frame") and shape.text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.font and para.font.size:
                        try:
                            size_pt = para.font.size.pt
                            result["font_sizes"].append(size_pt)
                        except Exception:
                            pass
                    for run in para.runs:
                        if run.font and run.font.name:
                            result["font_families"].add(run.font.name)

            # Extract chart info - safe check (GraphicFrame always has .chart attr but raises ValueError)
            try:
                from pptx.shapes.graphfrm import GraphicFrame

                if isinstance(shape, GraphicFrame) and shape.has_chart:
                    chart_info = extract_chart_colors(ppt_path, shape.chart)
                    result["charts_found"] += 1
                    for c in chart_info.get("palette", []):
                        if c not in result["palette"]:
                            result["palette"].append(c)
                    ct = chart_info.get("chart_type", "unknown")
                    result["chart_types"][ct] = result["chart_types"].get(ct, 0) + 1
                    if chart_info.get("has_data_labels"):
                        result["has_data_labels_total"] += 1
                    if chart_info.get("has_legend"):
                        result["has_legend_total"] += 1
            except Exception:
                pass

    # Convert sets to lists for JSON serialization
    result["font_families"] = sorted(result["font_families"])[:10]
    result["font_sizes"] = sorted(set(result["font_sizes"]))[:20]
    result["shapes_checked"] = shapes_checked
    result["total_slides"] = total_slides

    return result


def build_institution_dna(results: list[dict]) -> dict:
    """Aggregate per-institution style DNA from multiple PPT results."""
    from collections import defaultdict

    inst_data = defaultdict(
        lambda: {
            "files": [],
            "charts_found": 0,
            "palette": [],
            "font_families": [],
            "font_sizes": [],
            "chart_types": {},
            "data_label_ratio": 0,
            "legend_ratio": 0,
        }
    )

    for r in results:
        inst = r.get("institution")
        if not inst:
            continue

        d = inst_data[inst]
        d["files"].append(r["file"])
        d["charts_found"] += r.get("charts_found", 0)

        # Merge palette (keep order, deduplicate)
        for c in r.get("palette", []):
            if c not in d["palette"]:
                d["palette"].append(c)

        # Merge fonts
        for f in r.get("font_families", []):
            if f not in d["font_families"]:
                d["font_families"].append(f)

        # Merge font sizes (all of them)
        d["font_sizes"].extend(r.get("font_sizes", []))

        # Sum chart types
        for ct, count in r.get("chart_types", {}).items():
            d["chart_types"][ct] = d["chart_types"].get(ct, 0) + count

        d["data_label_ratio"] += r.get("has_data_labels_total", 0)
        d["legend_ratio"] += r.get("has_legend_total", 0)

    # Normalize
    output = {}
    for inst, d in inst_data.items():
        total_charts = max(d["charts_found"], 1)
        output[inst] = {
            "files_analyzed": d["files"],
            "charts_analyzed": d["charts_found"],
            "palette": d["palette"][:15],  # Top 15 colors
            "top_colors_used": d["palette"][:10],
            "font_families": d["font_families"][:5],
            "font_sizes_used": sorted(set(d["font_sizes"]))[:15],
            "chart_types": dict(sorted(d["chart_types"].items(), key=lambda x: -x[1])),
            "data_label_ratio": round(d["data_label_ratio"] / total_charts, 2),
            "legend_ratio": round(d["legend_ratio"] / total_charts, 2),
        }

    return output


def main():
    """Main entry: scan all PPTs, extract DNA, write JSON."""
    logging.basicConfig(level=logging.INFO, format="%(levelname)s | %(message)s")

    if not PPT_DIR.exists():
        logger.error(f"PPT directory not found: {PPT_DIR}")
        logger.info(f"Looking for: {TOOLKIT_DIR}")
        if TOOLKIT_DIR.exists():
            logger.info(f"Toolkit exists, contents: {[p.name for p in TOOLKIT_DIR.iterdir()]}")
        return

    ppt_files = sorted(PPT_DIR.glob("*.pptx"))  # .ppt not supported by python-pptx
    logger.info(f"Found {len(ppt_files)} PPT files in {PPT_DIR}")

    results = []
    for ppt_path in ppt_files:
        inst = identify_institution(ppt_path.stem)
        logger.info(f"Scanning: {ppt_path.name} (identified: {inst})")
        result = scan_ppt(ppt_path)
        results.append(result)
        logger.info(
            f"  -> {result.get('charts_found', 0)} charts, "
            f"{len(result.get('palette', []))} colors, "
            f"{len(result.get('font_families', []))} fonts"
        )

    # Build per-institution DNA
    dna = build_institution_dna(results)

    # Add summary
    output = {
        "version": "v54",
        "generated_at": "2026-07-30",
        "institutions_with_data": len(dna),
        "institutions": dna,
        "all_files": [r["file"] for r in results],
    }

    # Write output
    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    with open(OUTPUT_PATH, "w", encoding="utf-8") as f:
        json.dump(output, f, ensure_ascii=False, indent=2)

    logger.info(f"\nDone! Output: {OUTPUT_PATH}")
    logger.info(f"  PPTs scanned: {len(ppt_files)}")
    logger.info(f"  Institutions with data: {len(dna)}")
    for inst, data in dna.items():
        logger.info(f"  {inst}: {data['charts_analyzed']} charts, {len(data['palette'])} palette colors")

    return output


if __name__ == "__main__":
    main()
