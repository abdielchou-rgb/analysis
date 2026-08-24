"""V51 PPTX Exporter — 对标 505 九大投行图表 toolkit 的机构级 PPT。

从 templates/ 目录读取机构模板（.potx），渲染报告内容。

核心对标:
  505 九大投行图表 toolkit 中 McKinsey/GS/BCG/JP Morgan 的真实 PPT 布局:
  - McKinsey: 每页一个核心结论作标题, 图表在下方, 数据标签直接标在图上
  - Goldman Sachs: 标题+副标题, 左侧图表右侧注释, 配色偏冷
  - BCG: 2x2 矩阵 + 增长-份额框架
  - 罗兰贝格: 436页模板中的 2x2 矩阵布局
"""

from __future__ import annotations
import logging
import os
from pathlib import Path
from typing import Optional

logger = logging.getLogger("v51.export.pptx")

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    _HAS_PPTX = True
except ImportError:
    _HAS_PPTX = False

from utils.chart_config import get_palette

TEMPLATE_DIR = Path(__file__).resolve().parent.parent / "templates"

TEMPLATE_MAP = {
    "cicc": "cicc/slides.potx",
    "goldman_sachs": "gs/slides.potx",
    "morgan_stanley": "ms/slides.potx",
    "mckinsey": "mck/slides.potx",
    "bcg": "bcg/slides.potx",
    "citic": "citic/slides.potx",
    "academic": "academic/slides.potx",
    "jpmorgan": "jpm/slides.potx",
    "bain": "bain/slides.potx",
    "deloitte": "deloitte/slides.potx",
    "ey": "ey/slides.potx",
    "kpmg": "kpmg/slides.potx",
    "pwc": "pwc/slides.potx",
    "citi": "citi/slides.potx",
    "htsc": "htsc/slides.potx",
    "csc": "csc/slides.potx",
}


def export_pptx(report_md: str, style_id: str = "cicc",
                chart_paths: dict[str, str] = None,
                output_path: str = "") -> Optional[str]:
    """将报告导出为 PPTX。

    从 templates/{style_id}/slides.potx 读取模板，填充内容。

    Args:
        report_md: markdown 格式的报告正文
        style_id: 机构风格 ID
        chart_paths: {"bar": "path.png", "sensitivity": "path.png", ...}
        output_path: 输出路径，默认 autogenerate
    """
    if not _HAS_PPTX:
        # P1-4（2026-08-07）：原逻辑 return None 静默跳过，导致调用方无法区分
        # "PPTX 未产出"与"导出失败"。改为明确报错。
        raise RuntimeError(
            "python-pptx 未安装，无法生成 PPTX。请执行: pip install python-pptx"
        )

    style = get_palette(style_id)
    colors = style["palette"]

    # 尝试加载模板
    template_path = None
    if style_id in TEMPLATE_MAP:
        tmpl = TEMPLATE_DIR / TEMPLATE_MAP[style_id]
        if tmpl.exists():
            template_path = str(tmpl)

    try:
        prs = Presentation(template_path) if template_path else Presentation()
    except Exception:
        prs = Presentation()

    # 清除默认 slide（如果有模板则保留模板格式）
    for _ in range(len(prs.slides) - 1, 0, -1):
        try:
            rId = prs.slides._sldIdLst[-1].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
            if rId:
                prs.part.drop_rel(rId)
        except Exception:
            pass

    # 从 markdown 解析章节和段落
    lines = report_md.split('\n')
    slide_content = []
    current_title = ""

    for line in lines:
        if line.startswith('# ') or line.startswith('## '):
            if current_title:
                slide_content.append({"title": current_title, "body": [], "has_chart": False, "tables": []})
            level = line.count('#')
            current_title = line.lstrip('#').strip()
        elif line.strip() and current_title:
            slide_content[-1]["body"].append(line.strip()) if slide_content and current_title else None

        # 检测 markdown 表格（该 slide 包含表格）
        if '|' in line and '---' not in line and current_title and slide_content:
            slide_content[-1].setdefault("tables", []).append(line)

    if current_title and (not slide_content or slide_content[-1].get("title") != current_title):
        slide_content.append({"title": current_title, "body": [], "has_chart": False, "tables": []})

    # 检查报告中的图表引用并分配
    chart_list = []
    if chart_paths:
        chart_list = sorted(chart_paths.items())
    # R30 模块3（排版修复）：chart_paths 为空时自动发现 output/charts 下的图
    # 根因：pipeline_fingerprint chart_paths=0 → 柯力 PPTX 21图全缺
    if not chart_list:
        try:
            from pathlib import Path as _P
            _charts_dir = _P(__file__).resolve().parent.parent / "output" / "charts"
            if _charts_dir.exists():
                _pngs = sorted(_charts_dir.glob("*.png"))
                chart_list = [(p.stem, str(p)) for p in _pngs[:30]]
                if chart_list:
                    logger.info("[PPTX-AUTO-CHART] 自动发现 %d 张图表", len(chart_list))
        except Exception as e:
            logger.warning("[PPTX-AUTO-CHART] 自动发现失败: %s", e)

    # 在关键章节插入图表（核心判断、财务分析、估值分析等）
    # R30 增强：每个含图表关键词的章节分配一张图（从图池顺序取，允许复用）
    chart_keywords = ["核心判断", "核心分歧", "财务", "估值", "收入", "利润", "竞争",
                      "KPI", "用户指标", "广告", "AI", "电商", "概况", "格局",
                      "风险", "现金流", "资产负债", "盈利", "DCF", "敏感性"]
    assigned_indices = set()
    chart_pool_idx = 0
    for sec in slide_content:
        # 章节标题含关键词 → 分配一张图
        matched = False
        for kw in chart_keywords:
            if kw in sec["title"]:
                matched = True
                break
        if matched and chart_list and chart_pool_idx < len(chart_list):
            sec["has_chart"] = True
            # 记录该章节用哪张图（chart_pool_idx 递增）
            sec["_chart_idx"] = chart_pool_idx
            chart_pool_idx += 1
            # 图足够时每章节用新图；图不足时循环复用
            if chart_pool_idx >= len(chart_list):
                chart_pool_idx = 0
        else:
            sec["_chart_idx"] = -1

    # 生成 slide
    slide_layout = prs.slide_layouts[1] if prs.slide_layouts else None
    body_font_style = {"size": Pt(10), "name": "Arial"}
    title_font_style = {"size": Pt(16), "name": "Arial", "bold": True}

    # 图表分配器
    chart_items = chart_list  # R30: 复用自动发现的 chart_list
    chart_idx = 0

    for i, section in enumerate(slide_content):
        title = section["title"]
        body = section["body"]

        # 选择 layout：有图用 Two Content (3)，有表格用 Title and Content (1)
        tables = section.get("tables", [])
        # R30: 用 _chart_idx 定位图表（-1=无图）
        _ci = section.get("_chart_idx", -1)
        should_have_chart = _ci >= 0 and chart_items and _ci < len(chart_items)
        has_tables = len(tables) > 2

        if should_have_chart:
            try:
                slide = prs.slides.add_slide(prs.slide_layouts[3])  # Two Content
            except Exception:
                slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
        else:
            try:
                slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
            except Exception:
                slide = prs.slides.add_slide(prs.slide_layouts[6])

        # 标题
        title_placeholder = None
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 0:
                title_placeholder = ph
                break
        if title_placeholder:
            title_placeholder.text = title.replace('### ', '').replace('## ', '').replace('# ', '')
        else:
            # fallback: add textbox
            txbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
            txbox.text_frame.text = title.replace('### ', '').replace('## ', '').replace('# ', '')

        # 标题
        if slide.shapes.title:
            slide.shapes.title.text = title
            for run in slide.shapes.title.text_frame.paragraphs[0].runs:
                run.font.size = Pt(16)
                run.font.color.rgb = RGBColor(0x00, 0x33, 0x66)

        # 内容区域
        body_shape = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5))
        tf = body_shape.text_frame
        tf.word_wrap = True

        chart_inserted = False
        should_have_chart = section.get("has_chart", False) and chart_items

        # 如果有表格，把表格渲染成文字列表嵌入
        tables = section.get("tables", [])
        if tables:
            # 提取表头作为 slide 正文
            table_lines = []
            for tbl_line in tables[:8]:
                cells = [c.strip() for c in tbl_line.split('|') if c.strip()]
                if cells:
                    table_lines.append(' | '.join(cells[:4]))
            if table_lines:
                table_text = '\n'.join(table_lines[:6])
                has_table_text = True
            else:
                has_table_text = False
        else:
            has_table_text = False

        body_text_idx = 0
        max_body_lines = 6 if should_have_chart or has_table_text else 15
        for para_text in body[:max_body_lines]:
            if not para_text.strip():
                continue
            if body_text_idx == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = para_text[:150]
            p.font.size = Pt(10)
            p.space_after = Pt(4)
            body_text_idx += 1

        # 嵌入表格文本
        if has_table_text and not chart_inserted:
            if body_text_idx > 0:
                p = tf.add_paragraph()
                p.text = ''
                p.space_after = Pt(2)
            for tbl_line in table_lines[:5]:
                p = tf.add_paragraph()
                p.text = tbl_line[:120]
                p.font.size = Pt(8)
                p.space_after = Pt(1)
                body_text_idx += 1

        # 嵌入图表（R30: 用 _chart_idx 精确定位）
        if _ci >= 0 and _ci < len(chart_items):
            ctype, cpath = chart_items[_ci]
            if Path(cpath).exists():
                try:
                    slide.shapes.add_picture(cpath, Inches(4.8), Inches(1.5), Inches(4.5), Inches(3.5))
                    chart_inserted = True
                except Exception as e:
                    # P1-4（2026-08-07）：原 pass 静默丢弃图表插入失败，
                    # 改为 warning 便于排查缺失/损坏的图表文件。
                    logger.warning("PPTX chart insert failed for %s: %s", cpath, e)

        if not chart_inserted and chart_items and not should_have_chart and _ci < 0:
            # 兜底：无关键词匹配的 slide 尝试嵌图
            for chart_type, chart_path in chart_items:
                if chart_path and Path(chart_path).exists():
                    try:
                        slide.shapes.add_picture(chart_path, Inches(4.8), Inches(1.5), Inches(4.5), Inches(3.5))
                        chart_inserted = True
                        break
                    except Exception as e:
                        # P1-4（2026-08-07）：改为 warning 而非静默 continue
                        logger.warning("PPTX fallback chart insert failed for %s: %s", chart_path, e)
                        continue

    # 保存
    out_path = output_path or str(Path("outputs") / f"report_{style_id}.pptx")
    Path("outputs").mkdir(exist_ok=True)
    prs.save(out_path)
    logger.info(f"PPTX exported: {out_path}")
    return out_path
